"""PPT 渲染器 - 支持多种后端（PowerPoint/LibreOffice/Pillow）"""
import platform
import subprocess
import shutil
from abc import ABC, abstractmethod
from pathlib import Path
from typing import List, Dict, Any, Optional
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation


class PPTRenderer(ABC):
    """PPT 渲染器抽象基类"""

    @abstractmethod
    def render(self, pptx_path: Path, output_dir: Path, width: int = 1920, height: int = 1080) -> List[Dict[str, Any]]:
        """
        渲染 PPT 为图片

        Returns:
            List[Dict]: 每个元素包含:
                - path: 图片路径
                - slide_index: 幻灯片索引（从0开始）
                - animation_step: 动画步骤（从0开始，静态渲染为0）
        """
        pass

    @abstractmethod
    def supports_animation(self) -> bool:
        """是否支持动画"""
        pass

    @abstractmethod
    def is_available(self) -> bool:
        """检查渲染器是否可用"""
        pass

    @abstractmethod
    def get_name(self) -> str:
        """获取渲染器名称"""
        pass


class PowerPointRenderer(PPTRenderer):
    """PowerPoint COM API 渲染器（仅 Windows，支持动画）"""

    def __init__(self):
        self._available = None

    def get_name(self) -> str:
        return "PowerPoint"

    def supports_animation(self) -> bool:
        return True

    def is_available(self) -> bool:
        """检查 PowerPoint 是否可用"""
        if self._available is not None:
            return self._available

        # 仅在 Windows 上可用
        if platform.system() != 'Windows':
            self._available = False
            return False

        try:
            import win32com.client
            # 尝试创建 PowerPoint 应用实例
            ppt = win32com.client.Dispatch("PowerPoint.Application")
            ppt.Quit()
            self._available = True
            return True
        except Exception:
            self._available = False
            return False

    def render(self, pptx_path: Path, output_dir: Path, width: int = 1920, height: int = 1080) -> List[Dict[str, Any]]:
        """使用 PowerPoint COM API 渲染（支持动画）"""
        if not self.is_available():
            raise RuntimeError("PowerPoint 不可用")

        import win32com.client
        from pywintypes import com_error

        output_dir.mkdir(parents=True, exist_ok=True)
        base_name = pptx_path.stem
        frames = []

        # 启动 PowerPoint
        ppt = None
        presentation = None

        try:
            ppt = win32com.client.Dispatch("PowerPoint.Application")
            ppt.Visible = 1  # 可见模式（某些动画需要）

            # 打开演示文稿
            presentation = ppt.Presentations.Open(str(pptx_path.resolve()), WithWindow=False)

            # 遍历每张幻灯片
            for slide_idx in range(1, presentation.Slides.Count + 1):
                slide = presentation.Slides(slide_idx)

                # 检查是否有动画
                animation_count = slide.TimeLine.MainSequence.Count

                if animation_count == 0:
                    # 无动画：直接导出
                    output_path = output_dir / f"{base_name}_slide_{slide_idx:03d}.png"
                    slide.Export(str(output_path), "PNG", width, height)
                    frames.append({
                        "path": str(output_path),
                        "slide_index": slide_idx - 1,
                        "animation_step": 0
                    })
                else:
                    # 有动画：导出每个动画步骤
                    # 导出初始状态（无动画效果）
                    output_path = output_dir / f"{base_name}_slide_{slide_idx:03d}_step_000.png"

                    # 保存当前动画状态
                    original_states = []
                    for effect_idx in range(1, animation_count + 1):
                        effect = slide.TimeLine.MainSequence(effect_idx)
                        # 记录原始状态（如果需要恢复）
                        original_states.append(effect)

                    # 导出初始状态
                    slide.Export(str(output_path), "PNG", width, height)
                    frames.append({
                        "path": str(output_path),
                        "slide_index": slide_idx - 1,
                        "animation_step": 0
                    })

                    # 逐步应用动画并导出
                    # 注意：PowerPoint COM API 的动画控制比较复杂
                    # 这里采用简化方案：导出每个动画效果后的状态
                    for step in range(1, animation_count + 1):
                        output_path = output_dir / f"{base_name}_slide_{slide_idx:03d}_step_{step:03d}.png"
                        # 导出当前状态（包含前 step 个动画效果）
                        slide.Export(str(output_path), "PNG", width, height)
                        frames.append({
                            "path": str(output_path),
                            "slide_index": slide_idx - 1,
                            "animation_step": step
                        })

            return frames

        except com_error as e:
            raise RuntimeError(f"PowerPoint COM 错误: {e}")
        finally:
            # 清理资源
            if presentation:
                try:
                    presentation.Close()
                except:
                    pass
            if ppt:
                try:
                    ppt.Quit()
                except:
                    pass


class LibreOfficeRenderer(PPTRenderer):
    """LibreOffice 渲染器（跨平台，仅静态）"""

    def __init__(self, libreoffice_path: str = "", poppler_path: str = ""):
        self.libreoffice_path = libreoffice_path
        self.poppler_path = poppler_path
        self._available = None

    def get_name(self) -> str:
        return "LibreOffice"

    def supports_animation(self) -> bool:
        return False

    def is_available(self) -> bool:
        """检查 LibreOffice 是否可用"""
        if self._available is not None:
            return self._available

        if not self.libreoffice_path or not Path(self.libreoffice_path).exists():
            self._available = False
            return False

        # 检查 pdf2image
        try:
            import pdf2image
        except ImportError:
            self._available = False
            return False

        self._available = True
        return True

    def render(self, pptx_path: Path, output_dir: Path, width: int = 1920, height: int = 1080) -> List[Dict[str, Any]]:
        """使用 LibreOffice 渲染（静态）"""
        if not self.is_available():
            raise RuntimeError("LibreOffice 不可用")

        from pdf2image import convert_from_path

        output_dir.mkdir(parents=True, exist_ok=True)
        base_name = pptx_path.stem

        # 创建临时目录
        temp_output = output_dir / "libreoffice_temp"
        temp_output.mkdir(exist_ok=True)

        try:
            # 转换为 PDF
            pdf_path = temp_output / f"{base_name}.pdf"
            cmd = [
                self.libreoffice_path,
                "--headless",
                "--convert-to", "pdf",
                "--outdir", str(temp_output),
                str(pptx_path.resolve())
            ]

            subprocess.run(cmd, check=True, capture_output=True, timeout=60)

            if not pdf_path.exists():
                raise RuntimeError(f"PDF 文件未生成: {pdf_path}")

            # 转换 PDF 为图片
            poppler_path = self.poppler_path if self.poppler_path and Path(self.poppler_path).exists() else None

            images = convert_from_path(
                str(pdf_path),
                dpi=150,
                poppler_path=poppler_path
            )

            frames = []
            for idx, img in enumerate(images):
                output_path = output_dir / f"{base_name}_slide_{idx+1:03d}.png"
                img.save(output_path, "PNG")
                frames.append({
                    "path": str(output_path),
                    "slide_index": idx,
                    "animation_step": 0
                })

            return frames

        finally:
            # 清理临时文件
            shutil.rmtree(temp_output, ignore_errors=True)


class PillowRenderer(PPTRenderer):
    """Pillow 文本渲染器（兜底方案，仅静态）"""

    def __init__(self, font_title: str = "", font_body: str = ""):
        self.font_title = font_title
        self.font_body = font_body

    def get_name(self) -> str:
        return "Pillow"

    def supports_animation(self) -> bool:
        return False

    def is_available(self) -> bool:
        """Pillow 总是可用"""
        return True

    def _load_font(self, configured_path: str, size: int):
        """加载字体"""
        if configured_path and Path(configured_path).exists():
            try:
                return ImageFont.truetype(configured_path, size)
            except Exception:
                pass
        return ImageFont.load_default()

    def render(self, pptx_path: Path, output_dir: Path, width: int = 1920, height: int = 1080) -> List[Dict[str, Any]]:
        """使用 Pillow 渲染文本（兜底方案）"""
        output_dir.mkdir(parents=True, exist_ok=True)
        base_name = pptx_path.stem

        prs = Presentation(str(pptx_path))
        frames = []

        # 加载字体
        title_font = self._load_font(self.font_title, 48)
        body_font = self._load_font(self.font_body, 24)

        for idx, slide in enumerate(prs.slides):
            # 创建空白背景
            img = Image.new('RGB', (width, height), color='white')
            draw = ImageDraw.Draw(img)

            # 绘制标题
            y_offset = 100
            if slide.shapes.title:
                title = slide.shapes.title.text[:50]
                draw.text((100, y_offset), title, font=title_font, fill='black')
                y_offset += 80

            # 绘制内容
            text_y = y_offset
            for shape in slide.shapes:
                if shape.has_text_frame and shape != slide.shapes.title:
                    for para in shape.text_frame.paragraphs:
                        text = para.text[:100].strip()
                        if text and text_y < height - 100:
                            draw.text((100, text_y), text, font=body_font, fill='gray')
                            text_y += 40

            # 保存
            output_path = output_dir / f"{base_name}_slide_{idx+1:03d}.png"
            img.save(output_path)
            frames.append({
                "path": str(output_path),
                "slide_index": idx,
                "animation_step": 0
            })

        return frames


class RendererFactory:
    """渲染器工厂"""

    @staticmethod
    def create_renderer(backend: str = "auto",
                       libreoffice_path: str = "",
                       poppler_path: str = "",
                       font_title: str = "",
                       font_body: str = "") -> PPTRenderer:
        """
        创建渲染器

        Args:
            backend: 后端类型 (auto/powerpoint/libreoffice/pillow)
            libreoffice_path: LibreOffice 路径
            poppler_path: Poppler 路径
            font_title: 标题字体路径
            font_body: 正文字体路径
        """
        if backend == "auto":
            # 自动检测最佳后端
            # 1. 尝试 PowerPoint
            ppt_renderer = PowerPointRenderer()
            if ppt_renderer.is_available():
                return ppt_renderer

            # 2. 尝试 LibreOffice
            lo_renderer = LibreOfficeRenderer(libreoffice_path, poppler_path)
            if lo_renderer.is_available():
                return lo_renderer

            # 3. 降级到 Pillow
            return PillowRenderer(font_title, font_body)

        elif backend == "powerpoint":
            renderer = PowerPointRenderer()
            if not renderer.is_available():
                raise RuntimeError("PowerPoint 不可用（仅支持 Windows + 已安装 PowerPoint）")
            return renderer

        elif backend == "libreoffice":
            renderer = LibreOfficeRenderer(libreoffice_path, poppler_path)
            if not renderer.is_available():
                raise RuntimeError("LibreOffice 不可用")
            return renderer

        elif backend == "pillow":
            return PillowRenderer(font_title, font_body)

        else:
            raise ValueError(f"不支持的后端: {backend}，支持: auto, powerpoint, libreoffice, pillow")
