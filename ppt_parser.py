"""PPTX 解析模块 - 提取幻灯片内容"""
import os
import sys
import json
from pathlib import Path
from typing import List, Dict, Any
from pptx import Presentation


class PPTParser:
    """PPTX 文档解析器"""

    def __init__(self, input_dir: str, temp_dir: str):
        self.input_dir = Path(input_dir)
        self.temp_dir = Path(temp_dir)
        self.temp_dir.mkdir(parents=True, exist_ok=True)

    def find_pptx_files(self) -> List[Path]:
        """查找 input 目录下的所有 pptx 文件"""
        return [f for f in self.input_dir.glob("*.pptx") if not f.name.startswith("~$")]

    def parse_slide(self, slide, slide_index: int) -> Dict[str, Any]:
        """解析单张幻灯片内容"""
        content = {
            "index": slide_index,
            "title": "",
            "texts": [],
            "notes": ""
        }

        # 提取标题
        if slide.shapes.title:
            content["title"] = slide.shapes.title.text.strip()

        # 提取所有文本
        for shape in slide.shapes:
            self._extract_shape_text(shape, content)

        # 提取备注
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            content["notes"] = notes_text

        return content

    def _extract_shape_text(self, shape, content: Dict[str, Any]):
        """递归提取 shape 中的文本（支持普通文本、表格、组合）"""
        if shape.has_text_frame:
            # 普通文本框、标题等
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text and text != content["title"]:
                    content["texts"].append(text)
        elif shape.shape_type == 19:  # TABLE
            # 提取表格中的文本
            try:
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        text = cell.text.strip()
                        if text and text != content["title"]:
                            content["texts"].append(text)
            except Exception:
                pass  # 忽略表格解析错误
        elif shape.shape_type == 6:  # GROUP
            # 递归提取组合对象中的文本
            try:
                for sub_shape in shape.shapes:
                    self._extract_shape_text(sub_shape, content)
            except Exception:
                pass  # 忽略组合解析错误

    def parse(self, pptx_path: Path) -> Dict[str, Any]:
        """解析整个 PPTX 文件"""
        prs = Presentation(str(pptx_path))

        result = {
            "filename": pptx_path.name,
            "total_slides": len(prs.slides),
            "slides": []
        }

        for idx, slide in enumerate(prs.slides):
            slide_content = self.parse_slide(slide, idx)
            result["slides"].append(slide_content)

        return result

    def export_to_text(self, pptx_path: Path, output_path: Path):
        """将 PPTX 内容导出为纯文本（用于 LLM 分析）"""
        data = self.parse(pptx_path)

        lines = []
        lines.append(f"# {data['filename']}\n")
        lines.append(f"共 {data['total_slides']} 张幻灯片\n")
        lines.append("=" * 50 + "\n\n")

        for slide in data["slides"]:
            lines.append(f"## 第 {slide['index'] + 1} 页\n")
            if slide["title"]:
                lines.append(f"标题: {slide['title']}\n")
            if slide["texts"]:
                lines.append("内容:\n")
                for text in slide["texts"]:
                    lines.append(f"  - {text}\n")
            if slide["notes"]:
                lines.append(f"备注: {slide['notes']}\n")
            lines.append("\n" + "-" * 30 + "\n\n")

        output_path.write_text("".join(lines), encoding="utf-8")
        return output_path
